"""Build the 20-minute cdot talk as a .pptx."""
import os
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR

HERE = os.path.dirname(os.path.abspath(__file__))
# Generated figures and built decks are artefacts, so they live in the repo's
# ignored output/ directory. Run make_figure1.py and make_figs.py first.
OUTPUT = os.path.normpath(os.path.join(HERE, "..", "..", "output"))
FIGS = os.path.join(OUTPUT, "talk_figs")
OUT = os.environ.get("DECK_OUT", os.path.join(OUTPUT, "cdot_talk_20min.pptx"))

SURFACE = RGBColor(0xFC, 0xFC, 0xFB)
INK = RGBColor(0x0B, 0x0B, 0x0B)
INK2 = RGBColor(0x52, 0x51, 0x4E)
MUTED = RGBColor(0x8A, 0x88, 0x80)
BLUE = RGBColor(0x2A, 0x78, 0xD6)
ORANGE = RGBColor(0xEB, 0x68, 0x34)
RULE = RGBColor(0xE0, 0xDE, 0xD7)
CODEBG = RGBColor(0xF1, 0xEF, 0xEA)

FONT = "Calibri"
MONO = "Consolas"

W, H = Inches(13.333), Inches(7.5)
MARGIN = Inches(0.72)
BODY_W = W - 2 * MARGIN

prs = Presentation()
prs.slide_width, prs.slide_height = W, H
BLANK = prs.slide_layouts[6]


def set_bg(slide, colour=SURFACE):
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = colour


def textbox(slide, left, top, width, height, align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP):
    tb = slide.shapes.add_textbox(left, top, width, height)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = 0
    tf.vertical_anchor = anchor
    tf.paragraphs[0].alignment = align
    return tf


def run(par, text, size=18, colour=INK, bold=False, font=FONT, italic=False):
    r = par.add_run()
    r.text = text
    r.font.size = Pt(size)
    r.font.color.rgb = colour
    r.font.bold = bold
    r.font.italic = italic
    r.font.name = font
    return r


def rect(slide, left, top, width, height, colour):
    from pptx.enum.shapes import MSO_SHAPE
    sh = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    sh.fill.solid()
    sh.fill.fore_color.rgb = colour
    sh.line.fill.background()
    sh.shadow.inherit = False
    return sh


# Per-slide speaking budget in seconds, in slide order. Content slides only:
# the backup slides at the end are not part of the 20-minute run.
DURATIONS = [
    15,   # 1  title
    45,   # 2  the task
    60,   # 3  why it is hard
    45,   # 4  the Python HGVS stack
    50,   # 5  UTA limits
    45,   # 6  cdot in one slide
    45,   # 7  architecture
    35,   # 8  Shariant
    50,   # 9  how the data is built
    40,   # 10 JSON format
    40,   # 11 three ways in
    55,   # 12 R1 coverage
    60,   # 13 R2 ClinVar vs UTA
    45,   # 14 R2 full-scale ClinVar
    45,   # 15 R2 Shariant
    50,   # 16 R3 throughput
    30,   # 17 R3 at scale
    45,   # 18 R4 cleaning intro
    55,   # 19 R4 cleaning result
    45,   # 20 R4 repaired / not repaired examples
    45,   # 21 R4 residual ceiling
    45,   # 22 R5 fallback
    50,   # 23 R5 safety measurement
    65,   # 24 R5 safety gate
    35,   # 25 limitations
    35,   # 26 summary
]
_clock = [0]
_budget = list(DURATIONS)


def _mmss(secs):
    return f"{secs // 60}:{secs % 60:02d}"


def notes(slide, text):
    """Prepend the running clock to each content slide's notes."""
    import re
    body = re.sub(r"^\d+:\d\d[–-]\d+:\d\d \([^)]*\)\n\n", "", text)
    if _budget:
        d = _budget.pop(0)
        start, end = _clock[0], _clock[0] + d
        _clock[0] = end
        body = f"{_mmss(start)}–{_mmss(end)}  ({d} s)\n\n" + body
    slide.notes_slide.notes_text_frame.text = body


def footer(slide, n):
    tf = textbox(slide, MARGIN, H - Inches(0.52), BODY_W, Inches(0.3))
    p = tf.paragraphs[0]
    run(p, "cdot  ·  SACGF", size=10, colour=MUTED)
    tf2 = textbox(slide, MARGIN, H - Inches(0.52), BODY_W, Inches(0.3), align=PP_ALIGN.RIGHT)
    run(tf2.paragraphs[0], str(n), size=10, colour=MUTED)


_counter = [0]


def new_slide(with_footer=True):
    slide = prs.slides.add_slide(BLANK)
    set_bg(slide)
    _counter[0] += 1
    if with_footer:
        footer(slide, _counter[0])
    return slide


def heading(slide, title, kicker=None):
    top = Inches(0.55)
    if kicker:
        tf = textbox(slide, MARGIN, top, BODY_W, Inches(0.3))
        run(tf.paragraphs[0], kicker.upper(), size=12, colour=BLUE, bold=True)
        top = top + Inches(0.36)
    tf = textbox(slide, MARGIN, top, BODY_W, Inches(0.75))
    run(tf.paragraphs[0], title, size=31, colour=INK, bold=True)
    y = top + Inches(0.78)
    rect(slide, MARGIN, y, Inches(1.6), Pt(3), BLUE)
    return y + Inches(0.34)


def apply_bullet(par, char, colour, marL, indent, font="Arial"):
    """Real PowerPoint hanging bullet, so wrapped lines stay aligned."""
    from pptx.oxml import parse_xml
    from pptx.oxml.ns import nsdecls
    pPr = par._p.get_or_add_pPr()
    pPr.set("marL", str(int(marL)))
    pPr.set("indent", str(int(indent)))
    hexval = "%02X%02X%02X" % (colour[0], colour[1], colour[2])
    pPr.append(parse_xml(f'<a:buClr {nsdecls("a")}><a:srgbClr val="{hexval}"/></a:buClr>'))
    pPr.append(parse_xml(f'<a:buFont {nsdecls("a")} typeface="{font}"/>'))
    pPr.append(parse_xml(f'<a:buChar {nsdecls("a")} char="{char}"/>'))


def bullets(slide, items, top, width=None, size=18.5, gap=Pt(11)):
    """items: str, or (level, str), or ('code', str), or ('note', str)."""
    width = width or BODY_W
    tf = textbox(slide, MARGIN, top, width, H - top - Inches(0.7))
    first = True
    for item in items:
        level, text, kind = 0, item, "bullet"
        if isinstance(item, tuple):
            a, text = item
            if a in ("code", "note", "big"):
                kind = a
            else:
                level = a
        p = tf.paragraphs[0] if first else tf.add_paragraph()
        first = False
        p.space_before = Pt(2) if level else gap
        p.space_after = Pt(0)
        p.line_spacing = 1.06
        if kind == "code":
            p.space_before = Pt(10)
            run(p, text, size=size - 3.5, colour=BLUE, font=MONO)
        elif kind == "note":
            run(p, text, size=size - 3, colour=INK2, italic=True)
        elif kind == "big":
            run(p, text, size=size + 6, colour=INK, bold=True)
        elif level == 0:
            apply_bullet(p, "▪", BLUE, Inches(0.30), -Inches(0.30))
            run(p, text, size=size, colour=INK)
        else:
            apply_bullet(p, "–", MUTED, Inches(0.30) + Inches(0.36) * level,
                         -Inches(0.24))
            run(p, text, size=size - 2, colour=INK2)
    return tf


def picture(slide, name, top, height=None, width=None, max_w=None):
    """Fit the image inside (max_w x height) and centre it horizontally."""
    path = os.path.join(FIGS, name)
    max_w = max_w or BODY_W
    pic = slide.shapes.add_picture(path, 0, top)
    scale = min(max_w / pic.width, (height or pic.height) / pic.height)
    pic.width = int(pic.width * scale)
    pic.height = int(pic.height * scale)
    pic.left = int((W - pic.width) / 2)
    return pic


def example_col(slide, left, top, width, title, items, accent=BLUE, size=13.5):
    """A titled column of HGVS examples.

    items: (string, second_line, label). second_line is the repaired string when
    accent is BLUE, or the reason it cannot be repaired otherwise.
    """
    tf = textbox(slide, left, top, width, Inches(0.4))
    run(tf.paragraphs[0], title, size=17, colour=accent, bold=True)
    rect(slide, left, top + Inches(0.36), Inches(1.1), Pt(2.5), accent)
    tf = textbox(slide, left, top + Inches(0.62), width, H - top - Inches(1.3))
    first = True
    for broken, second, label in items:
        p = tf.paragraphs[0] if first else tf.add_paragraph()
        first = False
        p.space_before = Pt(0) if p is tf.paragraphs[0] else Pt(13)
        p.space_after = Pt(0)
        p.line_spacing = 1.05
        run(p, broken, size=size, colour=INK, font=MONO)
        p2 = tf.add_paragraph()
        p2.space_before = Pt(1)
        p2.space_after = Pt(0)
        p2.line_spacing = 1.05
        if accent is BLUE:
            run(p2, "→  ", size=size, colour=accent, font=MONO, bold=True)
            run(p2, second, size=size, colour=accent, font=MONO)
        else:
            run(p2, "✗  ", size=size, colour=accent, font=MONO, bold=True)
            run(p2, second, size=size - 0.5, colour=INK2)
        if label:
            run(p2, "   " + label, size=size - 1, colour=MUTED, italic=True)
    return tf


# ============================================================== 1. title
s = new_slide(with_footer=False)
rect(s, 0, 0, Inches(0.28), H, BLUE)
tf = textbox(s, Inches(1.3), Inches(2.15), Inches(11), Inches(1.3))
run(tf.paragraphs[0], "cdot", size=68, colour=INK, bold=True)
tf = textbox(s, Inches(1.3), Inches(3.3), Inches(10.6), Inches(1.6))
p = tf.paragraphs[0]
run(p, "Resolving more HGVS variants", size=28, colour=INK2)
p2 = tf.add_paragraph(); p2.space_before = Pt(6)
run(p2, "Fast, versioned RefSeq and Ensembl transcript data,\nstring repair, and safe version "
        "fallback", size=21, colour=MUTED)
rect(s, Inches(1.3), Inches(4.9), Inches(1.6), Pt(3), BLUE)
tf = textbox(s, Inches(1.3), Inches(5.3), Inches(10.6), Inches(1.2))
p = tf.paragraphs[0]
run(p, "Dave Lawrence", size=19, colour=INK, bold=True)
p2 = tf.add_paragraph(); p2.space_before = Pt(6)
run(p2, "SACGF  ·  github.com/SACGF/cdot  ·  cdotlib.org", size=16, colour=INK2)
notes(s, "0:00–0:30 (30 s)\n\nTitle. Add co-authors and affiliations before presenting.\n\n"
         "Opening line: the goal is to get as many real-world HGVS strings as possible to "
         "resolve to a genomic coordinate. cdot does that on three fronts: it supplies the "
         "versioned transcript-to-genome alignments the Python HGVS libraries need, it repairs "
         "malformed strings before they reach the parser, and it offers a safe fallback when "
         "the cited transcript version no longer exists.")

# ============================================================== 2. the task
s = new_slide()
top = heading(s, "The task: an HGVS string back to a genomic coordinate")
tf = textbox(s, MARGIN, top + Inches(0.1), BODY_W, Inches(1.4))
p = tf.paragraphs[0]
run(p, "NM_007294.4(BRCA1):c.68_69delAG", size=23, colour=INK, font=MONO, bold=True)
p2 = tf.add_paragraph(); p2.space_before = Pt(8)
run(p2, "↓   GRCh38", size=21, colour=BLUE, bold=True)
p3 = tf.add_paragraph(); p3.space_before = Pt(8)
run(p3, "NC_000017.11:g.43124030_43124031del", size=23, colour=BLUE, font=MONO, bold=True)
bullets(s, [
    "HGVS nomenclature is the international standard for describing sequence variants",
    "Clinical reports, ClinVar submissions, literature curation, lab importers, search boxes",
    "Almost every variant-curation pipeline has to do this conversion, at volume",
], top + Inches(2.0))
notes(s, "0:30–1:30 (1 min)\n\nSet up the problem in concrete terms. A c.HGVS is a position in a "
         "transcript. To do anything genomic with it (compare to a VCF, look it up in a "
         "population database, match it against another lab's classification) you have to put "
         "it back on the genome, as a g.HGVS on a genomic accession, and from there into "
         "whatever coordinate representation the pipeline uses.\n\nThe example is the BRCA1 "
         "185delAG founder variant.\n\nThis is routine, and it is done at volume: importers, "
         "search boxes, batch curation.")

# ============================================================== 3. why hard
s = new_slide()
top = heading(s, "Why it is harder than it looks")
bullets(s, [
    "The answer depends on a transcript version and a genome build",
    (1, "c.68 in NM_007294.3 need not be the same genomic base as in NM_007294.4"),
    (1, "you need that exact alignment: exon boundaries, CDS, and transcript-vs-genome indel gaps"),
    "Scale: ClinVar alone carries HGVS for over 3 million variants",
    "Input quality: Mutalyzer reported ~50% error rates in submitted descriptions over five years",
    "Source choice changes the answer: only 44% of putative loss-of-function variants were "
    "called LoF by both RefSeq and Ensembl annotation in ANNOVAR",
], top, size=19)
notes(s, "1:30–2:45 (1 min 15)\n\nThree separate difficulties.\n\n1. Version and build specific. "
         "The transcript accession names a sequence at a version, and that version is aligned to "
         "a build. Without the matching alignment you cannot convert at all.\n\n"
         "2. Volume, so it has to be fast.\n\n3. The strings are dirty, and the annotation source "
         "you choose is itself a decision that changes results. That 44% figure is the argument "
         "for covering both RefSeq and Ensembl rather than picking one.")

# ============================================================== 4. the stack
s = new_slide()
top = heading(s, "The Python HGVS stack", kicker="where cdot fits")
bullets(s, [
    "biocommons/hgvs reads transcripts through a pluggable data-provider interface",
    (1, "that interface has become a de facto standard: hgvs-weaver and others implement it too"),
    (1, "one backend written to it serves every client at once"),
    "PyHGVS loads transcripts from its own structures, and is no longer maintained",
    "cdot does not implement the HGVS grammar or the coordinate mapping. It is the data layer "
    "underneath, plus a thin repair layer in front of the parser.",
    ("code", "from cdot.hgvs.dataproviders import JSONDataProvider"),
    ("code", 'hdp = JSONDataProvider(["cdot.0.2.34.refseq.grch38.json.gz"])'),
    ("note", "That is the whole integration. Everything else in the pipeline is unchanged."),
], top, size=19)
notes(s, "2:45–3:45 (1 min)\n\nImportant framing point: cdot is not a competitor to "
         "biocommons/hgvs or to VariantValidator or Mutalyzer. It sits underneath them, as the "
         "provider of transcript coordinates.\n\nBe precise about the boundary: cdot does not "
         "own the HGVS grammar or the c-to-g mapping, those stay with biocommons/hgvs. But it "
         "is not purely passive data either: clean_hgvs() repairs the string before it reaches "
         "the parser, and the version fallback picks a transcript version the backend can "
         "actually serve. Those are the two ends of the pipeline that the parser cannot help "
         "with.\n\nBecause the data-provider interface is a de facto standard, improving the "
         "backend improves everything built on it.")

# ============================================================== 5. UTA limits
s = new_slide()
top = heading(s, "The standard backend is UTA", kicker="two separate problems")
bullets(s, [
    "It is a PostgreSQL database, so you either run it or reach it over the wire",
    (1, "the public server is remote: roughly 1 transcript per second, and many hospital and "
        "corporate networks block the port outright"),
    (1, "a local copy is fast, but it is a database to install, load and keep current: a real "
        "operational commitment for a group that just wants to resolve variants"),
    "And what is in it is limited: ~141,000 versioned alignments, RefSeq only",
    (1, "no Ensembl at all, and no T2T-CHM13v2.0"),
    (1, "limited historical depth: retired transcript versions drop out"),
], top, size=19)
notes(s, "3:45–5:00 (1 min 15)\n\nUTA is a good resource and this is not a criticism of it, it "
         "just has a shape that does not fit some use cases.\n\nMake the point that there are "
         "two independent problems here, and the deployment one bites first. You are offered a "
         "choice between slow-and-remote and fast-but-you-run-a-database. Neither is attractive "
         "if you are a diagnostic lab that wants to resolve a few million HGVS strings and not "
         "become a DBA. cdot's answer is a file: no server, no port, no schema to migrate.\n\n"
         "Only then the content: RefSeq only, no Ensembl, no T2T. The Ensembl gap matters in "
         "Europe and in research genomics generally.")

# ============================================================== 6. what cdot is
s = new_slide()
top = heading(s, "cdot in one slide")
bullets(s, [
    ("big", "1,954,742 versioned transcript/genome alignments"),
    "Every available RefSeq GFF3 and Ensembl GTF release, converted to gzipped JSON",
    "GRCh37, GRCh38 and T2T-CHM13v2.0; RefSeq and Ensembl together",
    "No database. Load one file, or call the REST API at cdotlib.org.",
    "Drop-in replacement for UTA in any biocommons/hgvs pipeline",
    "MANE Select and Ensembl canonical tags stored, so a gene symbol can pick a transcript",
    "MIT licence, pip install cdot",
], top, size=19)
notes(s, "5:00–6:00 (1 min)\n\nThe headline. Nearly 2 million alignments, both sources, three builds, "
         "no database.\n\nThe unit here is the transcript-version alignment: a particular "
         "transcript version aligned to a particular build. The same version counts separately "
         "per build, because that alignment is exactly what a resolution against that build "
         "needs.")

# ============================================================== 7. architecture
s = new_slide()
tf = textbox(s, MARGIN, Inches(0.42), BODY_W, Inches(0.45))
p = tf.paragraphs[0]
run(p, "How it fits together", size=24, colour=INK, bold=True)
run(p, "     build the data once (A), resolve variants against it (B)", size=15, colour=INK2)
picture(s, "figure1.png", Inches(0.98), height=Inches(5.9), max_w=Inches(12.4))
notes(s, "6:00–6:45 (45 s)\n\nOne slide for the whole architecture, then the rest of the talk "
         "walks through it.\n\nBand A is build time and happens once: the annotation releases go "
         "through the Snakemake pipeline and come out as gzipped JSON, which is either downloaded "
         "directly or loaded into Redis behind the REST API.\n\nBand B is what a client does. A "
         "possibly malformed HGVS string is cleaned, then biocommons/hgvs does the parsing and "
         "mapping, pulling transcript coordinates from one of the cdot providers and sequence from "
         "the sequence layer. The opt-in helpers on the left rewrite the input before any of that "
         "happens.\n\nBlue is cdot, grey is external. Zoomed versions of each band are in the "
         "backup slides if anyone wants the detail.")

# ============================================================== 7. Shariant
s = new_slide()
top = heading(s, "Where it came from: Shariant", kicker="motivation")
bullets(s, [
    "Australian Genomics Shariant pools variant classifications from clinical testing labs",
    "Submissions accumulate over many years and are continually updated",
    "Each one is written against whichever transcript version was current at the time",
    "Many of those versions have since been retired from the current annotation",
    ("big", "Resolve real historical strings, not just today's."),
], top, size=19)
notes(s, "6:00–6:45 (45 s)\n\nThis is why cdot exists and why it is shaped the way it is. A "
         "national variant-sharing platform inherits years of accumulated classifications. If "
         "you can only resolve current transcript versions, a large slice of that data is "
         "simply unusable.\n\nThat requirement drove two design decisions: ingest the complete "
         "release history, and be tolerant of malformed input.")

# ============================================================== 8. build
s = new_slide()
top = heading(s, "How the data is built")
bullets(s, [
    "Download the complete run of historical releases from the RefSeq and Ensembl FTP sites",
    (1, "RefSeq GFF3: 9 GRCh37, 22 GRCh38, 5 T2T releases (through RS_2025_08)"),
    (1, "plus NCBI's historical alignment file (RS_2024_08), for versions retired before "
        "the oldest archived release"),
    (1, "Ensembl GTF: 3 GRCh37, 36 GRCh38, 2 T2T (releases 81 to 116)"),
    (1, "UTA fills accessions that predate the oldest archived annotation files"),
    "Parse with HTSeq, normalise contig names via bioutils, derive CDS from start/stop codons",
    "Merge chronologically, newer release wins",
    (1, "not just recency: early RefSeq GFF3 files omit alignment-gap CIGARs that later ones carry"),
    "Snakemake pipeline, output is gzipped JSON",
], top, size=18)
notes(s, "6:45–8:00 (1 min 15)\n\nThe pipeline itself is not exotic. The interesting decision is "
         "going back through the FTP archive rather than taking the current annotation.\n\n"
         "The chronological merge is worth a sentence: letting newer data win is not only about "
         "recency, it also repairs gaps, because older annotation files were sometimes less "
         "complete than the later ones covering the same transcript version.")

# ============================================================== 9. format
s = new_slide()
top = heading(s, "The JSON format")
bullets(s, [
    "Per transcript: gene_name, hgnc, biotype, accession and version",
    "genome_builds is always a dict keyed by build, even in a single-build file",
    (1, "so one file can stand in for a GTF/GFF, or carry all three builds at once"),
    "Per build: contig, strand, cds_start, cds_end, and exons",
    ("code", "exon = [alt_start, alt_end, exon_id, cds_start, cds_end, gap]"),
    (1, "gap carries the transcript-vs-genome indel CIGAR, converted at query time"),
    "MANE Select, MANE Plus Clinical, RefSeq Select, Ensembl canonical tags",
    "schema_version at root, so a client rejects an incompatible file on load",
    ("note", "JSON because every language parses it fast with built-in libraries, and it "
             "serialises straight over HTTP. Same files drive the local provider and the REST API."),
], top, size=17.5)
notes(s, "8:00–9:00 (1 min)\n\nKeep this brisk, it is a reference slide.\n\nTwo points worth "
         "landing: the build is always a dict key, which is what lets one format serve both a "
         "single-build drop-in file and a combined multi-build response; and the gap field, "
         "which preserves transcript-vs-genome indels so downstream libraries can apply them "
         "during conversion.\n\nAlso: the per-release JSON files are published on the GitHub "
         "releases page, so they double as a fast-loading replacement for the corresponding "
         "GTF/GFF even if you never touch HGVS.")

# ============================================================== 10. access
s = new_slide()
top = heading(s, "Three ways to get at it")
bullets(s, [
    "Local JSON: JSONDataProvider loads a GRCh38 RefSeq file in ~11 s, then O(1) lookups",
    "REST: cdotlib.org, one request per transcript, or one batched prefetch() to warm the cache",
    "Ensembl TARK: EnsemblTarkDataProvider, for when you need Ensembl's own authoritative source",
    (1, "the only client we know of exposing TARK through the biocommons interface"),
    "Sequence comes from SeqRepo, cdot's FastaSeqFetcher (local genome FASTA), or both chained",
    (1, "FastaSeqFetcher splices exon ranges out of the genome: fully offline, no SeqRepo install"),
], top, size=18.5)
notes(s, "9:00–10:00 (1 min)\n\nThree access modes off the same data.\n\nThe sequence layer is a "
         "separate concern from the transcript layer and it trips people up: resolving a c.HGVS "
         "needs the transcript sequence as well as the alignment. SeqRepo is RefSeq-centric and "
         "does not hold every Ensembl transcript, which is why FastaSeqFetcher and "
         "ChainedSeqFetcher exist.")

# ============================================================== 11. R1 coverage
s = new_slide()
top = heading(s, "Coverage: 13.9× more alignments than UTA", kicker="result 1")
picture(s, "coverage.png", top, height=Inches(4.2), max_w=Inches(11.7))
tf = textbox(s, MARGIN, H - Inches(1.32), BODY_W, Inches(0.8))
p = tf.paragraphs[0]
run(p, "Historical depth (the full release run) plus Ensembl: 854,940 Ensembl accessions are "
       "absent from UTA entirely. T2T-CHM13v2.0 adds 186,644 alignments, a first for these "
       "libraries.", size=15, colour=INK2)
notes(s, "10:00–11:15 (1 min 15)\n\nTwo things drive the gain.\n\nFirst, historical depth. UTA "
         "does hold several versions per transcript, but cdot ingests the complete run of "
         "releases, so an NM_ version cited in a 2015 report still resolves after NCBI retired "
         "it.\n\nSecond, Ensembl, which UTA does not cover at all.\n\nAnd T2T: those alignments "
         "are absent from UTA and from every other backend biocommons/hgvs or PyHGVS can use. "
         "VEP can annotate against T2T, but VEP is a standalone annotation tool, not a "
         "transcript backend for these libraries.")

# ============================================================== 12. R2 clinvar
s = new_slide()
top = heading(s, "Does it resolve real variants?", kicker="result 2")
picture(s, "clinvar.png", top, height=Inches(4.0), max_w=Inches(11.7))
tf = textbox(s, MARGIN, H - Inches(1.5), BODY_W, Inches(1.0))
p = tf.paragraphs[0]
run(p, "Identical biocommons/hgvs engine, shared local SeqRepo, UTA release uta_20241220 loaded "
       "locally. Only the transcript-data layer differs.", size=15, colour=INK2)
p2 = tf.add_paragraph(); p2.space_before = Pt(6)
run(p2, "RefSeq is at parity. The whole gap is Ensembl, which UTA does not store.",
    size=16, colour=INK, bold=True)
notes(s, "11:15–12:30 (1 min 15)\n\nLeft: a seeded ClinVar sample of 2,818 descriptions spanning "
         "both sources. cdot 99.0%, UTA 70.3%. It is a sample rather than all of ClinVar because "
         "the comparison is gated by UTA's throughput.\n\nRight: where that gap comes from. On "
         "the committed RefSeq test pairs both backends hold every cited current version, so it "
         "is a dead heat. On Ensembl, cdot resolves 99.8% and UTA resolves nothing, because it "
         "stores no Ensembl alignments.\n\nBe honest about the framing: this is not cdot being "
         "cleverer on RefSeq, it is cdot covering a source UTA does not.")

# ============================================================== 13. R2 full scale
s = new_slide()
top = heading(s, "All of ClinVar, through cdot", kicker="result 2")
bullets(s, [
    ("big", "4,382,308 (g.HGVS, c.HGVS) pairs   →   99.4% resolved"),
    "Of the resolved variants, 99.96% reproduce ClinVar's own VCF coordinate exactly",
    "Scored as CHROM/POS/REF/ALT, not as a g.HGVS string",
    (1, "so 3'-shifting, del/dup spellings and tandem-repeat notation are not miscounted as errors"),
    "The 0.04% residual is representation and multi-mapping, not coordinate errors",
    (1, "1,116 indel / left-alignment differences; 406 SNVs in just 42 paralog and copy-number "
        "transcripts that map to more than one locus; 178 over-shuffled insertions"),
    ("note", "Caveat: ClinVar is dominated by a few large labs citing mostly current RefSeq "
             "versions. It is a clean public scale check, not an unbiased sample of what labs use."),
], top, size=18)
notes(s, "12:30–13:30 (1 min)\n\nThis is the scale check, and only cdot is in the loop, so every "
         "variant can be included.\n\nThe scoring choice matters and is worth defending: if you "
         "compare g.HGVS strings you punish yourself for equivalent representations. Comparing "
         "VCF coordinates against ClinVar's own VCF is the honest test.\n\nThe residual is "
         "mostly paralogs and copy-number genes with a constant per-transcript offset, plus "
         "indel normalisation differences.")

# ============================================================== 14. R2 shariant
s = new_slide()
top = heading(s, "On the historical clinical data that motivated it", kicker="result 2  ·  tier 2")
bullets(s, [
    "36,010 unique HGVS descriptions imported into Shariant over many years",
    (1, "1,654 of the 3,731 distinct transcripts are cited at more than one version"),
    ("big", "cdot 99.6%     UTA 93.7%"),
    "Of the strings cdot resolved and UTA could not (5.9% of the corpus):",
    (1, "99% were RefSeq versions UTA holds no GRCh38 alignment for"),
    (1, "1% were Ensembl transcripts, absent from UTA entirely"),
    ("note", "No ground-truth coordinate exists for this corpus, so the metric is resolution "
             "rate, not correctness. Aggregate constants from private production data."),
], top, size=18.5)
notes(s, "13:30–14:30 (1 min)\n\nClinVar showed the pipeline works at scale on current versions. "
         "This shows why the historical depth is not a curiosity.\n\nThis is the traffic a "
         "working clinical lab actually generates: variants classified years ago against "
         "versions that have since been retired.\n\nFlag the tier-2 limitation openly: private "
         "corpus, aggregate frozen constants, a referee cannot rerun it, and there is no ground "
         "truth so we report resolution rate only.")

# ============================================================== 15. R3 throughput
s = new_slide()
top = heading(s, "Throughput", kicker="result 3")
picture(s, "throughput.png", top, height=Inches(4.1), max_w=Inches(11.7))
tf = textbox(s, MARGIN, H - Inches(1.42), BODY_W, Inches(0.95))
p = tf.paragraphs[0]
run(p, "Sequence layer held constant (shared local SeqRepo), identical engine, so only the "
       "transcript-data layer varies.", size=15, colour=INK2)
p2 = tf.add_paragraph(); p2.space_before = Pt(6)
run(p2, "cdot's local data layer is roughly 30× a local UTA. One prefetch() call gives REST "
        "local-file speed.", size=16, colour=INK, bold=True)
notes(s, "14:30–15:45 (1 min 15)\n\nNote the log scale.\n\nThe public remote UTA at 0.1 per "
         "second is the number most people actually experience if they do not self-host.\n\n"
         "REST at 39/s is one HTTP request per transcript version. Batch those into a single "
         "prefetch() and the per-request overhead amortises across the whole set, the cache "
         "warms, and later lookups are in-memory hits, so it matches local JSON.\n\n"
         "Once the data layer is local the bottleneck is sequence fetching, not transcript "
         "lookup, which is why the two cdot modes converge.")

# ============================================================== 16. R3 at scale
s = new_slide()
top = heading(s, "The same at scale, in one process", kicker="result 3")
bullets(s, [
    "3,660,452 unique ClinVar (g.HGVS, c.HGVS) pairs, single local-JSON process",
    (1, "~92 minutes at 665 HGVS/s"),
    "REST matched it over the network after one cache-warming pass",
    (1, "21,277 distinct transcripts warmed in ~6 s, then ~83 minutes at 731 HGVS/s"),
    (1, "once warm, both providers resolve from an in-memory dict with no further network I/O"),
    "The same exhaustive pass against the public remote UTA extrapolates to hundreds of days",
], top, size=19)
notes(s, "15:45–16:30 (45 s)\n\nQuick slide. The point is that the benchmark numbers hold up "
         "over millions of variants rather than a microbenchmark, and that the REST provider is "
         "genuinely usable for bulk work once you warm it.\n\nThe difference between 92 and 83 "
         "minutes is within run-to-run variance, plausibly because the warmed REST cache holds "
         "only the few thousand transcripts touched, whereas local JSON holds everything.")

# ============================================================== 17. cleaning intro
s = new_slide()
top = heading(s, "The strings people actually type", kicker="result 4")
bullets(s, [
    "They arrive from report PDFs, spreadsheets, literature, and free-text search boxes",
    "Whitespace and non-printables from Word and PDF, lost casing, transposed punctuation, "
    "gene symbol and accession swapped, a trailing protein annotation",
    "The intended variant is obvious to a human reader. A strict grammar rejects it outright.",
    ("big", "The usual response is to report it invalid. cdot tries to repair it."),
    "clean_hgvs() is a pure string operation: no genome build, no sequence, no parser, "
    "no data provider",
    (1, "so it runs anywhere, as a pre-pass ahead of either HGVS library"),
    "Ordered pipeline of single-purpose fixes; every change returned as an HGVSFix record",
    (1, "code, message, before/after, severity, so a caller can log, audit or reject it"),
], top, size=18)
notes(s, "16:30–17:30 (1 min)\n\nThis is the part of cdot that is not about transcript data at "
         "all, and it is often what people end up using first.\n\nEmphasise the design "
         "constraint: no parser, no provider, no build. That is what lets it sit in a search box "
         "or an importer with no infrastructure.\n\nAnd emphasise that nothing is silent. Every "
         "repair comes back as a structured record, so user input is never quietly mutated.")

# ============================================================== 18. cleaning result
s = new_slide()
top = heading(s, "91.5% → 96.6% parseable on a real query stream",
              kicker="result 4  ·  tier 2")
picture(s, "cleaning.png", top, height=Inches(4.0), max_w=Inches(11.7))
tf = textbox(s, MARGIN, H - Inches(1.42), BODY_W, Inches(0.95))
p = tf.paragraphs[0]
run(p, "32,752 real queries typed into the HGVS search box of production VariantGrid-based "
       "platforms.", size=15, colour=INK2)
p2 = tf.add_paragraph(); p2.space_before = Pt(6)
run(p2, "+5.1% absolute · 1,678 strings rescued · zero regressions: no already-valid string "
        "was broken.", size=16, colour=INK, bold=True)
notes(s, "17:30–18:45 (1 min 15)\n\nThis is messy human input, not synthetic errors. Users treat "
         "the search box as a shortcut to jump to a variant, so the strings carry the damage of "
         "however they got there.\n\nThe rescues are dominated by the dullest possible problems: "
         "whitespace and casing. That is the honest finding, and it is good news, because those "
         "are exactly the fixes that are unambiguous and safe to apply automatically.\n\n"
         "Categories overlap, so the counts sum to more than 1,678.\n\nThe zero-regression claim "
         "is separately backed by a reproducible injection benchmark in the repo, which injects "
         "each fix category into clean ClinVar strings and confirms recovery with no "
         "regressions.")

# ============================================================== 18b. examples
s = new_slide()
top = heading(s, "What gets repaired, and what does not", kicker="result 4")
COL_W = Inches(6.0)
example_col(s, MARGIN, top, COL_W, "Repaired", [
    ("NM_000059.4: c.1A>G", "NM_000059.4:c.1A>G", "whitespace"),
    ("NM_000059.4:c.1delg", "NM_000059.4:c.1delG", "base casing"),
    ("NM_000059..4:c.1A>G", "NM_000059.4:c.1A>G", "punctuation"),
    ("BRCA2(NM_000059.4):c.1A>G", "NM_000059.4(BRCA2):c.1A>G", "wrapper swapped"),
    ("NM_000059.4:c.1A>G p.(Met1?)", "NM_000059.4:c.1A>G", "protein suffix"),
], accent=BLUE)
example_col(s, MARGIN + Inches(6.5), top, COL_W, "Still fails", [
    ("NM_000059.4:c.68_69", "no edit: the string was cut off", "25%"),
    ("c.68_69delAG", "no transcript was ever given", "25%"),
    ("000059.4:c.68del", "NM_ prefix dropped", "15%"),
    ("NM_000059.4:c.68_69ins5", "inserted bases never supplied", "3%"),
    ("NM_000059.4:c.(67+1_68-1)_(70+1_71-1)del", "legal HGVS, grammar gap", "7%"),
], accent=ORANGE)
tf = textbox(s, MARGIN, H - Inches(1.0), BODY_W, Inches(0.5))
run(tf.paragraphs[0], "Examples synthesised from public BRCA2 NM_000059.4, not real queries. "
                      "Percentages are the share of the 1,118 queries that still fail after "
                      "cleaning.", size=13, colour=MUTED, italic=True)
notes(s, "18:45–19:30 (45 s)\n\nMake it concrete. Left column: the whole rescue is unglamorous. "
         "A space after the colon, a lower-case base, a doubled dot, the gene and accession the "
         "wrong way round, a protein annotation pasted on the end. Every one of these is "
         "unambiguous, which is why it is safe to fix automatically.\n\nRight column: the failures "
         "are a different kind of thing. The first two, half the residual, are not errors we could "
         "be cleverer about, the information is simply absent. You cannot invent the edit the user "
         "did not type, or the transcript they never named. The dropped NM_ prefix is the frontier "
         "case, that one is fixable and is on the list. The ins5 has a position and a length but "
         "no bases. And the last is our own grammar's limitation, not the user's mistake: that is "
         "valid HGVS for an uncertain-range deletion that biocommons will not parse.\n\nSo the "
         "ceiling is not a cleaning-effort problem, it is mostly a missing-information problem.")

# ============================================================== 19. residual
s = new_slide()
top = heading(s, "The ceiling of string repair", kicker="result 4  ·  tier 2")
picture(s, "residual.png", top, height=Inches(3.85), max_w=Inches(11.7))
tf = textbox(s, MARGIN, H - Inches(1.5), BODY_W, Inches(1.0))
p = tf.paragraphs[0]
run(p, "1,118 queries still fail after cleaning. Most of what remains is out of scope for any "
       "string-level repair: the information was never supplied.", size=15, colour=INK2)
p2 = tf.add_paragraph(); p2.space_before = Pt(6)
run(p2, "Classified single-label under a fixed decision tree by an LLM: single-rater, so no "
        "inter-rater agreement is reported.", size=14, colour=MUTED, italic=True)
notes(s, "18:45–19:45 (1 min)\n\nWorth being candid here. Cleaning has a hard ceiling and this "
         "quantifies it.\n\nAbout half is truncated input or a bare variant body with no "
         "reference sequence. No string repair can invent a transcript the user never typed.\n\n"
         "About a third is genuinely fixable and marks the frontier for future cleaning "
         "rules.\n\nThe remainder splits evenly between a grammar gap, where the input is legal "
         "HGVS the biocommons grammar rejects, and pasted URLs or prose that should never have "
         "been parsed.\n\nState the classification limitation plainly if asked.")

# ============================================================== 20. R5 fallback
s = new_slide()
top = heading(s, "When the cited version no longer exists", kicker="result 5")
bullets(s, [
    "Retired transcript versions are routine in historical clinical data",
    "get_best_transcript_version() offers the nearest available version",
    (1, "up-then-down, closest, or latest, as the caller configures"),
    "Never applied automatically: it is a heuristic, returned as an HGVSFix the caller can reject",
    "This is a client-layer feature, not a backend one",
    (1, "biocommons/hgvs has no version fallback with any provider, so a UTA-backed pipeline "
        "gains nothing here even though UTA stores several versions"),
    ("big", "But when is the swap safe?"),
    "Only if the cited position lands on the same genomic base under the substitute",
    (1, "a report cites NM_007294.3:c.68, we hold only NM_007294.4. Same genomic base?"),
    (1, "a new version can move it: re-annotated CDS, changed exons, a new alignment gap"),
], top, size=18, gap=Pt(9))
notes(s, "19:45–20:45 (1 min)\n\nThis is the feature that closes the loop on the Shariant "
         "motivation.\n\nThe strategy question (which version do I pick) is easy. The hard "
         "question is the safety one: substituting a different version can silently move the "
         "variant, and a silently moved clinical variant is much worse than a failed "
         "lookup.\n\nSo the default is to refuse, and every substitution is reported.")

# ============================================================== 21. R5 safety
s = new_slide()
top = heading(s, "How often is a version bump actually safe?", kicker="result 5")
picture(s, "version_safety.png", top, height=Inches(3.35), max_w=Inches(11.7))
tf = textbox(s, MARGIN, H - Inches(1.85), BODY_W, Inches(1.3))
p = tf.paragraphs[0]
run(p, "Measured directly on released cdot data: because cdot stores the alignment rather than "
       "the sequence, a bump can only move a coordinate if it changes that alignment.",
    size=15, colour=INK2)
p2 = tf.add_paragraph(); p2.space_before = Pt(6)
run(p2, "Partial drift (some variants move, some do not) is rare: 0.7% of RefSeq bumps, 0.4% "
        "Ensembl. Only 5.5% of RefSeq accessions ever drift at all.", size=15, colour=INK2)
notes(s, "20:45–21:45 (1 min)\n\nTwo ways to read the same data. Per version bump: does this "
         "bump preserve every coding coordinate. Per coding base: what is the chance a random "
         "variant is unaffected, which is closer to what a user experiences.\n\nBoth sources sit "
         "at 96% or better per bump. RefSeq is slightly lower than Ensembl here because the "
         "historical GRCh38 alignments bring in much older versions, which drift more.\n\n"
         "Two useful properties. First, the dangerous case, a partial bump "
         "that moves some variants but not others, is rare. When a coordinate does move, it is "
         "usually the whole CDS, driven by a re-annotation of the coding region. Second, drift "
         "is concentrated in a small set of accessions, so the risky transcripts form a short "
         "list you can flag rather than treating every bump as unsafe.")

# ============================================================== 22. R5 gate
s = new_slide()
top = heading(s, "The safety gate cdot ships", kicker="result 5")
bullets(s, [
    "Comparing the two versions' genomic coordinates does not work",
    (1, "the requested version is retired, and is usually not aligned to the build you are in"),
    "So compare something build-independent: the intrinsic CDS structure",
    (1, "CDS length plus the coding-exon segment lengths, in transcript coordinates"),
    (1, "verified identical across GRCh37 / GRCh38 / T2T for 99.3% of RefSeq, 100% of Ensembl"),
    (1, "so read the requested version's structure from whichever build still carries it, the "
        "substitute's from the target build, and swap only if the two match"),
    "A structure change stands in for a coordinate move: 99.9% of drifting RefSeq bumps and "
    "100% of drifting Ensembl bumps change it, and 100% of structure-unchanged bumps are "
    "genomically preserved",
    "Refinements: CDS alignment gaps must match, and for a cited UTR position, UTR length too",
    ("big", "Validated on ClinVar: of 4,024,794 accepted substitutions, 2 moved a coordinate."),
    (1, "both one transcript re-aligned to a different locus, now shipped in a small blocklist"),
    ("code", "is_version_substitution_safe()   ·   intrinsic_cds_structure()"),
], top, size=16.5, gap=Pt(8))
notes(s, "21:45–23:15 (1 min 30)\n\nThe key idea, and probably the most interesting engineering "
         "in the paper. Spend the time here.\n\nYou want to know whether swapping version N for "
         "version M moves the variant. The obvious check compares the two versions' genomic "
         "coordinates, but that fails exactly when you need it most: the version you were asked "
         "for is retired, which is usually why it is not in the build you are working in, so "
         "there is no coordinate to compare against.\n\nThe way out is that the check does not "
         "have to happen in the build you are resolving in. The intrinsic CDS structure is a "
         "property of the transcript sequence, not of the genome, so it is the same number "
         "wherever that version appears. We measured that: identical across GRCh37, GRCh38 and "
         "T2T for 99.3% of RefSeq and 100% of Ensembl versions. That is what makes the whole "
         "thing work, because cdot carries every build. The retired version may be gone from "
         "GRCh38 but still present in the GRCh37 or T2T data, and its structure read there is "
         "just as good.\n\nIt also "
         "catches transient reverts, a coordinate that goes A to B back to A, which a "
         "genomic-bracket check between neighbours cannot see.\n\nThe two failures out of four "
         "million were a single transcript re-placed to another locus, which the structure check "
         "cannot see by construction, so cdot compares genomic CDS maps when both are loaded and "
         "ships a blocklist for when they are not.")

# ============================================================== 23. limitations
s = new_slide()
top = heading(s, "Limitations")
bullets(s, [
    "Resolution is only as current as the ingested releases",
    (1, "a version published after the last ingested release needs the dataset regenerated; "
        "automating that is the next planned improvement"),
    "The production-corpus results are aggregate constants from private data, not "
    "referee-reproducible",
    "The residual error taxonomy is single-rater LLM classification, no κ reported",
    "FastaSeqFetcher reproduces an Ensembl transcript exactly, but a curated RefSeq sequence "
    "can differ from the genome at a few positions",
    "Canonical transcript selection is a sensible default for search, not a clinical decision",
], top, size=18.5)
notes(s, "23:15–24:00 (45 s)\n\nRun through these quickly and without hedging. Being explicit "
         "about the tier-2 data and the single-rater taxonomy tends to defuse the obvious "
         "reviewer questions rather than invite them.")

# ============================================================== 24. summary
s = new_slide()
top = heading(s, "Summary")
bullets(s, [
    "The goal is to resolve as many real-world HGVS strings as possible, correctly",
    "cdot supplies the versioned transcript/genome alignments the Python HGVS libraries need",
    "1,954,742 alignments, 13.9× UTA: RefSeq and Ensembl, GRCh37, GRCh38 and T2T-CHM13v2.0",
    "Roughly 30× a local UTA on the identical engine, with no database to run",
    "clean_hgvs() rescued 5.1% more of a real production query stream, with zero regressions",
    "Version fallback behind a safety gate validated on 4 million ClinVar substitutions",
    ("code", "pip install cdot     github.com/SACGF/cdot     cdotlib.org     MIT"),
], top, size=19)
notes(s, "24:00–24:45 (45 s), then questions.\n\nClose on the framing you opened with: existing "
         "pipelines gain coverage and speed without changing how they work. Two lines of code.\n\n"
         "Likely questions:\n"
         "· How do you keep it current? Regenerating from new releases; automation is planned.\n"
         "· How big are the files? Gzipped JSON, GRCh38 RefSeq loads in about 11 seconds.\n"
         "· Relationship to VariantValidator / Mutalyzer? Complementary. They validate and "
         "normalise descriptions; cdot supplies the transcript-coordinate layer.\n"
         "· Why not just use VEP for T2T? VEP is a standalone annotation tool, not a transcript "
         "backend for these libraries.")


# ============================================================== backup: band A
s = new_slide()
top = heading(s, "Data generation, in detail", kicker="backup  ·  figure 1, band a")
picture(s, "figure1_band_a.png", top + Inches(0.85), height=Inches(3.0), max_w=Inches(12.4))
notes(s, "Backup slide. Band A of Figure 1 at full width.\n\nUse if someone asks about the build "
         "pipeline: which sources go in, what the Snakemake pipeline does to them, and what comes "
         "out. Note that UTA appears here as an input, filling alignments that were never published "
         "in any GFF3 or GTF release.")

# ============================================================== backup: band B
s = new_slide()
top = heading(s, "Variant resolution, in detail", kicker="backup  ·  figure 1, band b")
picture(s, "figure1_band_b.png", top + Inches(0.1), height=Inches(4.3), max_w=Inches(12.4))
notes(s, "Backup slide. Band B of Figure 1 at full width.\n\nUse if someone asks how the pieces "
         "fit at call time: the three data providers behind the one biocommons interface, the "
         "separate sequence layer, and where clean_hgvs() and the opt-in helpers sit relative to "
         "the parser.")

prs.save(OUT)
print("talk length:", _mmss(_clock[0]))
print("wrote", OUT, "-", len(prs.slides.__iter__.__self__._sldIdLst), "slides")
